"""
文档解析服务
支持多种文档格式的文本提取
"""

import io
from typing import List, Dict, Optional
import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup


class DocumentChunk:
    """文档切片"""
    def __init__(
        self,
        text: str,
        page: Optional[int] = None,
        heading: Optional[str] = None,
        metadata: Optional[Dict] = None
    ):
        self.text = text
        self.page = page
        self.heading = heading
        self.metadata = metadata or {}


class DocumentParser:
    """文档解析器"""
    
    @staticmethod
    def parse_pdf(file_content: bytes) -> List[DocumentChunk]:
        """解析 PDF 文件"""
        chunks = []
        
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    chunks.append(DocumentChunk(
                        text=text.strip(),
                        page=page_num,
                        metadata={"page_width": page.width, "page_height": page.height}
                    ))
        
        return chunks
    
    @staticmethod
    def parse_docx(file_content: bytes) -> List[DocumentChunk]:
        """解析 DOCX 文件"""
        chunks = []
        doc = DocxDocument(io.BytesIO(file_content))
        
        current_heading = None
        current_text = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # 检测标题
            if para.style.name.startswith('Heading'):
                if current_text:
                    chunks.append(DocumentChunk(
                        text="\n".join(current_text),
                        heading=current_heading
                    ))
                    current_text = []
                current_heading = text
            else:
                current_text.append(text)
        
        # 添加最后一个块
        if current_text:
            chunks.append(DocumentChunk(
                text="\n".join(current_text),
                heading=current_heading
            ))
        
        return chunks
    
    @staticmethod
    def parse_txt(file_content: bytes) -> List[DocumentChunk]:
        """解析 TXT 文件"""
        text = file_content.decode('utf-8', errors='ignore')
        return [DocumentChunk(text=text.strip())]
    
    @staticmethod
    def parse_html(file_content: bytes) -> List[DocumentChunk]:
        """解析 HTML 文件"""
        soup = BeautifulSoup(file_content, 'html.parser')
        
        # 移除 script 和 style 标签
        for script in soup(["script", "style"]):
            script.decompose()
        
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks_text = '\n'.join(line for line in lines if line)
        
        return [DocumentChunk(text=chunks_text)]
    
    @staticmethod
    def parse_xlsx(file_content: bytes) -> List[DocumentChunk]:
        """解析 XLSX 文件"""
        chunks = []
        wb = load_workbook(io.BytesIO(file_content), read_only=True)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows_text = []
            
            for row in sheet.iter_rows(values_only=True):
                row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    rows_text.append(row_text)
            
            if rows_text:
                chunks.append(DocumentChunk(
                    text='\n'.join(rows_text),
                    heading=f"工作表: {sheet_name}"
                ))
        
        return chunks
    
    @staticmethod
    def parse_pptx(file_content: bytes) -> List[DocumentChunk]:
        """解析 PPTX 文件"""
        chunks = []
        prs = Presentation(io.BytesIO(file_content))
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            title = None
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title
                        title = shape.text.strip()
                    else:
                        slide_text.append(shape.text.strip())
            
            if slide_text:
                chunks.append(DocumentChunk(
                    text='\n'.join(slide_text),
                    page=slide_num,
                    heading=title or f"幻灯片 {slide_num}"
                ))
        
        return chunks
    
    @staticmethod
    def parse_md(file_content: bytes) -> List[DocumentChunk]:
        """解析 Markdown 文件"""
        text = file_content.decode('utf-8', errors='ignore')
        return [DocumentChunk(text=text.strip())]
    
    @classmethod
    def parse(cls, file_content: bytes, file_type: str) -> List[DocumentChunk]:
        """根据文件类型解析文档"""
        parsers = {
            'pdf': cls.parse_pdf,
            'docx': cls.parse_docx,
            'txt': cls.parse_txt,
            'html': cls.parse_html,
            'xlsx': cls.parse_xlsx,
            'pptx': cls.parse_pptx,
            'md': cls.parse_md,
        }
        
        parser = parsers.get(file_type.lower())
        if not parser:
            raise ValueError(f"不支持的文件类型: {file_type}")
        
        return parser(file_content)

